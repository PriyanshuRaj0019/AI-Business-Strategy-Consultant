from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def generate_ppt(company_name: str, report: str) -> str:
    """
    Generate a PowerPoint presentation from the AI report.
    """

    reports_dir = Path("reports")
    reports_dir.mkdir(exist_ok=True)

    ppt_path = reports_dir / f"{company_name}_Business_Report.pptx"

    prs = Presentation()

    # -----------------------
    # Title Slide
    # -----------------------
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = f"{company_name}"
    slide.placeholders[1].text = "AI Business Strategy Consultant"

    # -----------------------
    # Content Slides
    # -----------------------

    sections = report.split("##")

    for section in sections:

        section = section.strip()

        if not section:
            continue

        lines = section.split("\n")

        title = lines[0]

        body = "\n".join(lines[1:])

        layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = title

        textbox = slide.placeholders[1].text_frame

        textbox.clear()

        p = textbox.paragraphs[0]
        p.text = body[:3500]
        p.font.size = Pt(18)

    prs.save(str(ppt_path))

    return str(ppt_path)